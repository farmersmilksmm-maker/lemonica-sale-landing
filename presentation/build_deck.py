#!/usr/bin/env python3
"""LEMONICA — Hollywood district deck, 2 slides, light Lemonica palette."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

IVORY = RGBColor(0xF7, 0xF3, 0xE7)      # фон
CARD = RGBColor(0xFF, 0xFF, 0xFF)       # карточки
INK = RGBColor(0x26, 0x25, 0x1B)        # основной текст
MUTED = RGBColor(0x6E, 0x6A, 0x5C)      # вторичный текст
LEMON = RGBColor(0xE3, 0xB6, 0x2C)      # акцент-форма
LEMON_DEEP = RGBColor(0x8C, 0x6A, 0x0A) # акцент-текст (контраст ~4,5:1 на айвори)
GREEN = RGBColor(0x5C, 0x6B, 0x4C)      # второй акцент
LINE = RGBColor(0xE2, 0xDC, 0xC8)       # рамки
BAR_DIM = RGBColor(0xCB, 0xC4, 0xA8)    # столбец 2025

SERIF = "Georgia"
SANS = "Arial"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = IVORY
    bg.line.fill.background(); bg.shadow.inherit = False
    return s


def tx(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for (text, size, color, bold, italic, font) in para:
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.italic = italic; r.font.name = font
    return box


PDF_LINK = "https://farmersmilksmm-maker.github.io/lemonica-sale-landing/PDF090226_CoStar_YoungCircle.pdf"


def footer(s, page):
    box = s.shapes.add_textbox(Inches(0.6), Inches(7.04), Inches(11.6), Inches(0.4))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Источник: CoStar Group, отчёт по 1801–1848 N Young Cir, 02.09.2026 · PDF: "
    r1.font.size = Pt(9); r1.font.color.rgb = MUTED; r1.font.name = SANS
    r2 = p.add_run()
    r2.text = "PDF090226_CoStar_YoungCircle.pdf"
    r2.font.size = Pt(13); r2.font.color.rgb = LEMON_DEEP; r2.font.name = SANS
    r2.font.bold = True; r2.font.underline = True
    r2.hyperlink.address = PDF_LINK
    tx(s, 12.45, 7.04, 0.4, 0.35, [[(str(page), 10, MUTED, True, False, SANS)]], align=PP_ALIGN.RIGHT)


LOGO = "/Users/denysharbuzov/Documents/Limonika_Sale_Landing_2026-09-02/presentation/lemonica_logo.png"


def lemon_mark(s, x, y, size=0.9):
    """Настоящий логотип LEMONICA (прозрачный PNG из заготовок, без вырезок)."""
    s.shapes.add_picture(LOGO, Inches(x), Inches(y), Inches(size), Inches(size))


def header(s, kicker, title_runs, sub=None, title_size=38):
    tx(s, 0.6, 0.45, 9.0, 0.3, [[(kicker, 12, GREEN, True, False, SANS)]])
    tx(s, 0.6, 0.82, 12.1, 1.1, [[(t, title_size, c, b, i, f) for (t, c, b, i, f) in
                               [(r[0], r[2], r[3], r[4], r[5]) for r in title_runs]]], line_spacing=1.02)
    if sub:
        tx(s, 0.6, 1.95, 12.1, 0.45, [[(sub, 14, MUTED, False, False, SANS)]])


def stat_card(s, x, y, w, h, big, label, sub=None, accent=LEMON_DEEP):
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.adjustments[0] = 0.06
    card.fill.solid(); card.fill.fore_color.rgb = CARD
    card.line.color.rgb = LINE; card.line.width = Pt(1)
    card.shadow.inherit = False
    # лимонная метка слева
    tab = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + 0.28), Inches(0.07), Inches(h - 0.56))
    tab.fill.solid(); tab.fill.fore_color.rgb = accent; tab.line.fill.background(); tab.shadow.inherit = False
    pad = 0.3
    paras = [[(big, 27, accent, True, False, SERIF)],
             [(label, 11.5, INK, True, False, SANS)]]
    if sub:
        paras.append([(sub, 10.5, MUTED, False, False, SANS)])
    box = tx(s, x + pad, y + pad - 0.02, w - 2 * pad - 0.1, h - 2 * pad, paras, line_spacing=1.0)
    ps = box.text_frame.paragraphs
    ps[0].space_after = Pt(8)
    if sub:
        ps[2].space_before = Pt(4)
        ps[1].space_after = Pt(4)


# ============ SLIDE 1 — рост и люди (10 миль) ============
s = slide()
lemon_mark(s, 11.85, 0.28)
header(s, "ГОЛЛИВУД · ФЛОРИДА · 01",
       [("Район растёт: ", 38, INK, True, False, SERIF),
        ("+90 000 жителей за 5 лет", 38, LEMON_DEEP, True, False, SERIF)],
       "Радиус 10 миль от ресторана — прогноз CoStar, 2025–2030")

# chart: two bars
cx, cy, ch = 0.9, 2.9, 3.05
v25, v30 = 1207030, 1297766
vmax = 1360000.0
bw = 1.45
for i, (label, val, color, txt_color) in enumerate([
        ("2025", v25, BAR_DIM, MUTED),
        ("2030 (прогноз)", v30, LEMON, LEMON_DEEP)]):
    gx = cx + 0.8 + i * 3.1
    bh = (val / vmax) * (ch - 0.7)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(gx), Inches(cy + ch - 0.7 - bh), Inches(bw), Inches(bh))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background(); bar.shadow.inherit = False
    tx(s, gx - 0.8, cy + ch - 0.7 - bh - 0.44, bw + 1.6, 0.4,
       [[(f"{val/1e6:.2f} млн".replace(".", ","), 21, txt_color, True, False, SERIF)]], align=PP_ALIGN.CENTER)
    tx(s, gx - 0.8, cy + ch - 0.36, bw + 1.6, 0.3, [[(label, 12.5, INK, True, False, SANS)]], align=PP_ALIGN.CENTER)
# контекст переписи 2020 (город) — честная сноска, не кольцо
tx(s, cx + 0.1, cy + ch + 0.06, 7.2, 0.25, [[
    ("Контекст: перепись 2020 — город Голливуд 153 067 жителей; CoStar даёт оценки 2025/2030 для радиусных колец.",
     9, MUTED, False, True, SANS)]])
# pill with growth — под диаграммой, не пересекает подписи столбцов
pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx + 0.35), Inches(cy + ch + 0.42), Inches(5.6), Inches(0.52))
pill.adjustments[0] = 0.5
pill.fill.solid(); pill.fill.fore_color.rgb = LEMON; pill.line.fill.background(); pill.shadow.inherit = False
tx(s, cx + 0.35, cy + ch + 0.5, 5.6, 0.36, [[("+7,5% за 5 лет · +90 736 человек", 14, RGBColor(0x3D, 0x2E, 0x04), True, False, SANS)]], align=PP_ALIGN.CENTER)

# right: 4 cards 2x2
stat_card(s, 8.0, 2.72, 2.45, 1.68, "$71 363", "доход домохозяйств", "медианный · 10 миль")
stat_card(s, 10.55, 2.72, 2.45, 1.68, "41 год", "средний возраст", "радиус 10 миль")
stat_card(s, 8.0, 4.56, 2.45, 1.68, "$451 881", "цена жилья", "медианная · 10 миль")
stat_card(s, 10.55, 4.56, 2.45, 1.68, "463 898", "домохозяйств", "+7,6% к 2030")
footer(s, 1)

# ============ SLIDE 2 — стройка и спрос ============
s = slide()
lemon_mark(s, 11.85, 0.28)
header(s, "ГОЛЛИВУД · ФЛОРИДА · 02",
       [("Город застраивается: ", 38, INK, True, False, SERIF),
        ("спрос опережает предложение", 38, LEMON_DEEP, True, False, SERIF)],
       "Субмаркет Hollywood / Dania Beach и рынок Fort Lauderdale — CoStar, Q3 2026", title_size=34)

stat_card(s, 0.7, 2.75, 3.9, 1.85, "868", "квартир построено за 12 месяцев", "в субмаркете Hollywood / Dania Beach")
stat_card(s, 4.72, 2.75, 3.9, 1.85, "624", "квартир строится прямо сейчас", "+2,9% к жилому фонду района")
stat_card(s, 8.74, 2.75, 3.9, 1.85, "+1,1%", "рост аренды за год", "2-й результат из 9 субмаркетов")

stat_card(s, 0.7, 4.8, 3.9, 1.85, "6 345", "квартир в стройке по рынку", "инвестиции $1,5 млрд за год")
stat_card(s, 4.72, 4.8, 3.9, 1.85, "920 000", "рабочих мест в рынке", "поток гостей на ланч и ужин")
stat_card(s, 8.74, 4.8, 3.9, 1.85, "311", "квартир — новый дом у Young Circle", "комплекс Radius, N Young Cir")
footer(s, 2)

OUT = "/Users/denysharbuzov/Documents/Limonika_Sale_Landing_2026-09-02/presentation/Hollywood_District_CoStar_2slides.pptx"
prs.save(OUT)
print("saved", OUT)
